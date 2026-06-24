"""Fill slide 3 with 'How to run queries' navigation info."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def box(slide, text, left, top, w, h,
        fill=None, font_size=9, bold=False, color="1A1A2E",
        align=PP_ALIGN.LEFT, italic=False, font_name=None, wrap=True):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    if font_name: run.font.name = font_name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    return shape

NAVY="1A1A2E"; SPARQL_BLUE="216BC2"; SQL_SLATE="37474F"
WHITE="FFFFFF"; LIGHT_BG="F8F9FA"; MID_BG="EAECEF"
CODE_BG="1C2333"; CODE_FG="D4D4D4"
DARK_TEXT="2C2C2C"; MID_TEXT="555555"

prs = Presentation("bgg_agents_slides.pptx")
slide = prs.slides[2]

# Remove existing placeholder shapes
sp_tree = slide.shapes._spTree
for elem in list(sp_tree):
    tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
    if tag in ('sp', 'pic', 'graphicFrame', 'cxnSp', 'grpSp'):
        sp_tree.remove(elem)

# ── Title bar ─────────────────────────────────────────────────────
box(slide, "  How to Run Queries",
    0, 0, 13.33, 0.72, fill=NAVY, font_size=24, bold=True, color=WHITE)

# ── Section A label ───────────────────────────────────────────────
box(slide, "  A  Comparing the four KG agents",
    0.15, 0.84, 6.30, 0.40, fill=SPARQL_BLUE, font_size=13, bold=True, color=WHITE)

box(slide, "Run any cell in these Jupyter notebooks:",
    0.15, 1.30, 6.30, 0.26, font_size=9, color=MID_TEXT)

# Notebook table
NB_ROWS = [
    ("bgg_sparql_qa.ipynb",      "SPARQL agent",       "GraphDB (localhost:7200)"),
    ("bgg_graphrag_qa.ipynb",    "GraphRAG agent",     "ChromaDB vector store"),
    ("bgg_longcontext_qa.ipynb", "Long-context agent", "Full text to Claude Sonnet"),
    ("bgg_gemini_qa.ipynb",      "Gemini agent",       "Full text to Gemini 2.5 Flash"),
]
NB_FILLS = ["232B3E", "1C2333", "232B3E", "1C2333"]

for i, (fname, agent, backend) in enumerate(NB_ROWS):
    y = 1.60 + i * 0.48
    # filename
    box(slide, f"  {fname}", 0.15, y, 3.20, 0.36,
        fill=NB_FILLS[i], font_size=8.5, bold=True, color="7EC8E3", font_name="Consolas", wrap=False)
    # agent name
    box(slide, agent, 3.40, y, 1.60, 0.36,
        fill=NB_FILLS[i], font_size=8.5, bold=True, color=WHITE)
    # backend
    box(slide, backend, 5.05, y, 1.40, 0.36,
        fill=NB_FILLS[i], font_size=8, color="AAAAAA")

# Tip
box(slide,
    "Each notebook has a golden test section (A-G questions) at the bottom. "
    "Run any ask(...) cell to try a live query.",
    0.15, 3.58, 6.30, 0.36, fill="EBF5FB", font_size=8.5, color=DARK_TEXT)

# ── Divider ───────────────────────────────────────────────────────
# thin vertical rule between the two halves
rule = slide.shapes.add_textbox(Inches(6.60), Inches(0.84), Inches(0.03), Inches(5.60))
rule.fill.solid()
rule.fill.fore_color.rgb = rgb("CCCCCC")

# ── Section B label ───────────────────────────────────────────────
box(slide, "  B  Comparing SPARQL vs SQL (multi-hop)",
    6.75, 0.84, 6.40, 0.40, fill=SQL_SLATE, font_size=13, bold=True, color=WHITE)

box(slide, "Run these Python scripts from the terminal:",
    6.75, 1.30, 6.40, 0.26, font_size=9, color=MID_TEXT)

# Script table
SC_ROWS = [
    ("bgg_sparql_compare.py", "Claude generates + executes SPARQL",  "GraphDB (must be running)"),
    ("bgg_sql_compare.py",    "Claude generates + executes SQL",      "bgg.db  (SQLite, local)"),
]
SC_FILLS = ["232B3E", "1C2333"]

for i, (fname, desc, backend) in enumerate(SC_ROWS):
    y = 1.60 + i * 0.56
    box(slide, f"  {fname}", 6.75, y, 3.40, 0.36,
        fill=SC_FILLS[i], font_size=8.5, bold=True, color="7EC8E3", font_name="Consolas", wrap=False)
    box(slide, desc, 10.22, y, 2.90, 0.36,
        fill=SC_FILLS[i], font_size=8, color=WHITE)
    box(slide, backend, 6.75, y+0.36, 6.37, 0.20,
        fill=SC_FILLS[i], font_size=7.5, italic=True, color="AAAAAA")

# Run commands
box(slide, "Run commands:", 6.75, 2.90, 6.40, 0.24, font_size=9, bold=True, color=DARK_TEXT)
for i, cmd in enumerate([
    "uv run bgg_sparql_compare.py",
    "uv run bgg_sql_compare.py",
]):
    box(slide, f"  {cmd}", 6.75, 3.18 + i*0.36, 6.40, 0.30,
        fill=CODE_BG, font_size=9, color=CODE_FG, font_name="Consolas", wrap=False)

box(slide,
    "Both scripts loop through QSQL1-8 and print results. "
    "GraphDB must be running locally on port 7200 for the SPARQL script.",
    6.75, 3.96, 6.40, 0.36, fill="EBF5FB", font_size=8.5, color=DARK_TEXT)

# ── Bottom note ───────────────────────────────────────────────────
box(slide,
    "Prerequisites: .env file with ANTHROPIC_API_KEY (and GEMINI_API_KEY for the Gemini notebook).  "
    "GraphDB repo 'bgg' must be loaded with bgg_main.ttl for SPARQL queries.",
    0.15, 6.88, 13.00, 0.40, fill="FFF9E6", font_size=8, color="5D4037")

prs.save("bgg_agents_slides.pptx")
print("Slide 3 filled. Done.")
