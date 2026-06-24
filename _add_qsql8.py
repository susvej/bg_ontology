"""Add QSQL8 (anchorless family size at any depth) everywhere it belongs."""
import sys, json
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl

# ── Shared helpers ────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def box(slide, text, left, top, w, h,
        fill=None, font_size=9, bold=False, color="1A1A2E",
        align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name=None):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    if font_name: run.font.name = font_name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    return shape

def code_block(slide, code_text, left, top, w, h):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("1C2333")
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(8)
        run.font.name = "Consolas"
        run.font.color.rgb = rgb("D4D4D4")
    return shape

NAVY="1A1A2E"; SPARQL_BLUE="216BC2"; SQL_SLATE="37474F"
WIN_GREEN="1E8C69"; BUG_RED="C0392B"; WHITE="FFFFFF"; ORANGE="D46A08"

# ════════════════════════════════════════════════════════════════════
#  1. REBUILD SLIDE 12 with 8 rows
# ════════════════════════════════════════════════════════════════════
prs = Presentation("bgg_agents_slides.pptx")
slide12 = prs.slides[11]

# Keep title bar only
title_elem = None
for shape in slide12.shapes:
    try:
        fc = str(shape.fill.fore_color.rgb).upper()
        if fc == "1A1A2E" and shape.top == 0:
            title_elem = shape._element; break
    except: pass

sp_tree = slide12.shapes._spTree
for elem in list(sp_tree):
    tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
    if tag in ('sp','pic','graphicFrame','cxnSp','grpSp'):
        if elem is not title_elem:
            sp_tree.remove(elem)

# Column positions (unchanged)
CQ=0.15; WQ=3.20; CSL=3.40; WSL=1.85; CSS=5.30; WSS=1.85
CQL=7.20; WQL=1.75; CQS=9.00; WQS=1.85

HDR_Y=0.82; HDR_H=0.36
box(slide12,"Question",CQ,HDR_Y,WQ,HDR_H,fill=NAVY,font_size=9,bold=True,color=WHITE)
box(slide12,"SPARQL lines",CSL,HDR_Y,WSL,HDR_H,fill=SPARQL_BLUE,font_size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
box(slide12,"SPARQL score",CSS,HDR_Y,WSS,HDR_H,fill=SPARQL_BLUE,font_size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
box(slide12,"SQL lines",CQL,HDR_Y,WQL,HDR_H,fill=SQL_SLATE,font_size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
box(slide12,"SQL score",CQS,HDR_Y,WQS,HDR_H,fill=SQL_SLATE,font_size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

ROWS = [
    ("QSQL1",
     "Sea theme + route/network mechanic,\nnot owned by any of the 3 players",
     "Find me a game about the sea — whether it appears as a theme or a category, "
     "and whether the label used is 'Nautical', 'maritime', 'naval', or 'sailing' — "
     "that also features a route-building or network-building mechanic, and that none "
     "of the three of us (Susanne Vejdemo, Gertrude Young, and Maja Brown) already own. "
     "Order by rating and show the top results.",
     7,10,15,10,False),
    ("QSQL2",
     "Games by designer of Gertrude's 8+ games,\nnot owned by Susanne",
     "Suggest highly-rated games that I, Susanne Vejdemo, do not own, but that were "
     "designed by the same person who designed a game that Gertrude Young owns and "
     "rates 8 or above.",
     12,10,10,10,False),
    ("QSQL3",
     "Multi-constraint: designer + mechanic\ncross-player filter",
     "Find a highly-rated game that neither Susanne Vejdemo nor Gertrude Young owns, "
     "where the game was designed by someone who also designed a game that Susanne rates "
     "8 or higher, and where the game also features a mechanic found in a game that "
     "Gertrude rates 8 or higher.",
     20,10,25,10,False),
    ("QSQL4",
     "Susanne's nearest game-neighbors\n(most shared games)",
     "Who are Susanne Vejdemo's closest 'game neighbors' in this community? Find the "
     "other players who share the most games in common with Susanne, and show how many "
     "games they share. Order by overlap descending.",
     7,9,9,9,False),
    ("QSQL5",
     "Recommendation via 3-hop designer chain",
     "Recommend a highly-rated game (8 or above) that neither Susanne Vejdemo nor "
     "Gertrude Young owns, found through a three-hop designer chain: start from a game "
     "that Gertrude rates 8 or higher, follow the designer to another game they made, "
     "then follow a co-designer on that second game to find the final recommendation. "
     "Show the chain: Gertrude's game, the linking designer, the bridge game, the "
     "co-designer, and the recommended game.",
     25,10,40,10,False),
    ("QSQL6",
     "Transitive community (all players reachable\nthrough shared ownership)",
     "Who is in Susanne Vejdemo's gaming community? Find ALL players reachable from "
     "Susanne through chains of shared game ownership — where two players are 'connected' "
     "if they share at least one game in common, and the community is the full transitive "
     "closure of this connection relationship (i.e. it includes Susanne's direct "
     "game-neighbors, their game-neighbors, and so on to any depth). List all community "
     "members by name and report the total count.",
     4,10,15,7,True),
    ("QSQL7",
     "Ticket to Ride family\n(recursive reimplementation traversal)",
     "Show me all games in the Ticket to Ride game family — starting from the original "
     "Ticket to Ride (BGG ID 9209), find every game that reimplements it, every game "
     "that reimplements those, and so on recursively to any depth. List each game's "
     "name and rating, ordered by rating descending.",
     5,10,8,10,False),
    ("QSQL8",
     "Total reimplementation family size\n(any depth, no anchor game)",
     "Which games have the largest total reimplementation family — counting all games "
     "that have reimplemented them, directly or through any number of reimplementation "
     "generations? Return every game that has been reimplemented at least once, with "
     "its full family size, ordered largest first.",
     7,10,12,10,False),
]

XA_H=0.36; GAP_AB=0.02; GAP_PAIR=0.02
XB_H={"QSQL1":0.40,"QSQL2":0.22,"QSQL3":0.28,"QSQL4":0.24,
      "QSQL5":0.42,"QSQL6":0.44,"QSQL7":0.28,"QSQL8":0.22}

y = HDR_Y + HDR_H + 0.04
for i,(qid,abbrev,full_q,slines,sscore,qlines,qscore,is_key) in enumerate(ROWS):
    row_fill = "E3F2FD" if is_key else ("FAFAFA" if i%2==0 else "F0F0F0")
    xb_fill  = "EAF3FB" if is_key else ("F3F3F3" if i%2==0 else "E8E8E8")

    box(slide12,f"{qid}  {abbrev}",CQ,y,WQ,XA_H,fill=row_fill,font_size=8.5,bold=is_key)
    box(slide12,str(slines),CSL,y,WSL,XA_H,fill=row_fill,font_size=10,bold=is_key,align=PP_ALIGN.CENTER)
    box(slide12,f"{sscore}/10",CSS,y,WSS,XA_H,
        fill=(WIN_GREEN if sscore==10 else ORANGE),font_size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    box(slide12,str(qlines),CQL,y,WQL,XA_H,fill=row_fill,font_size=10,bold=is_key,align=PP_ALIGN.CENTER)
    box(slide12,f"{qscore}/10",CQS,y,WQS,XA_H,
        fill=(WIN_GREEN if qscore==10 else BUG_RED),font_size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

    yb = y + XA_H + GAP_AB
    box(slide12,full_q,CQ+0.30,yb,WQ+(CQS+WQS-CQ-0.30)-(CQ+0.30)+CQ,XB_H[qid],
        fill=xb_fill,font_size=7,bold=False,italic=True,color="444444")
    y = yb + XB_H[qid] + GAP_PAIR

box(slide12,
    "  Key: QSQL6 SQL bug returns 203 instead of 202.  QSQL6, QSQL7, QSQL8 all show "
    "SPARQL property paths vs SQL WITH RECURSIVE — QSQL8 is the cleanest anchorless example.",
    0.15,y+0.04,13.00,0.40,fill="FFF3CD",font_size=9,color="5D4037")

print(f"Slide 12 rebuilt — content ends at y={y:.2f} in")

# ════════════════════════════════════════════════════════════════════
#  2. ADD QSQL8 CODE SLIDE  (insert before existing slide 15 = key takeaways)
#     python-pptx can't insert mid-deck, so append then rely on presenter order
# ════════════════════════════════════════════════════════════════════
layout = prs.slide_layouts[6]
s_q8 = prs.slides.add_slide(layout)

# Title
box(s_q8,"  QSQL8 — Reimplementation Family Size (Any Depth, No Anchor)",
    0,0,13.33,0.72,fill=NAVY,font_size=20,bold=True,color=WHITE)

box(s_q8,
    "Which games have the largest total reimplementation family — at any depth, across the whole database?",
    0.15,0.80,13.00,0.38,font_size=9.5,italic=True,color="444444")

SPARQL_Q8 = """\
PREFIX bgg: <https://raw.githubusercontent.com/susvej/bg_ontology/>

SELECT ?gameName (COUNT(DISTINCT ?descendant) AS ?familySize) WHERE {
  ?game bgg:hasName ?gameName .
  ?descendant bgg:reimplements+ ?game .
}
GROUP BY ?game ?gameName
ORDER BY DESC(?familySize)"""

SQL_Q8 = """\
WITH RECURSIVE descendants(orig_id, desc_id) AS (
  SELECT older_id, newer_id
  FROM game_reimplements

  UNION

  SELECT d.orig_id, gr.newer_id
  FROM descendants d
  JOIN game_reimplements gr
    ON gr.older_id = d.desc_id
)
SELECT g.name, COUNT(*) AS family_size
FROM descendants d
JOIN games g ON g.id = d.orig_id
GROUP BY g.id, g.name
ORDER BY family_size DESC;"""

def section_label(slide, text, left, top, w, fill_color):
    box(slide,"  "+text,left,top,w,0.38,fill=fill_color,font_size=13,bold=True,color=WHITE)

section_label(s_q8,"SPARQL  (7 lines)",0.15,1.26,6.30,SPARQL_BLUE)
code_block(s_q8,SPARQL_Q8,0.15,1.68,6.30,2.70)

section_label(s_q8,"SQL  (12 lines + WITH RECURSIVE)",6.85,1.26,6.30,SQL_SLATE)
code_block(s_q8,SQL_Q8,6.85,1.68,6.30,2.70)

# Highlight the key token
box(s_q8,"bgg:reimplements+",
    0.60,2.52,4.20,0.36,fill="1A3A5C",font_size=10,bold=True,color="7EC8E3",
    font_name="Consolas")
box(s_q8,"change from depth-1 query: add one  +  token. SQL must restructure entirely.",
    0.60,2.90,5.50,0.26,font_size=8.5,color=SPARQL_BLUE)

# Results comparison
box(s_q8,"Both return 861 games. Top results:",
    0.15,4.54,13.00,0.26,font_size=9,bold=True,color="333333")

RESULTS = [
    ("CATAN","33"),("Star Wars: Epic Duels","25"),("Love Letter","25"),
    ("Ticket to Ride","21"),("Escape Room: The Game","21"),
    ("Carcassonne","16"),("Zombicide","15"),
]
for j,(name,n) in enumerate(RESULTS):
    col = j % 4
    row = j // 4
    box(s_q8,f"  {n:>3}  {name}",
        0.15 + col*3.28, 4.84 + row*0.30, 3.20, 0.28,
        fill=("1C2333" if j%2==0 else "22293F"),
        font_size=8,color="D4D4D4",font_name="Consolas")

# Evaluation callout
box(s_q8,
    "SPARQL advantage: the depth-1 query (bgg:reimplements without +) and this any-depth query differ "
    "by a single character. SQL requires adding a full WITH RECURSIVE block to compute the transitive "
    "closure before any aggregation is possible — a fundamentally different query structure. "
    "Results confirmed identical between SPARQL and SQL (both verified against live data).",
    0.15,5.60,13.00,0.60,fill="F0F4FF",font_size=8.5,color="333333")

print("QSQL8 code slide added")

# ════════════════════════════════════════════════════════════════════
#  3. UPDATE report.xlsx
# ════════════════════════════════════════════════════════════════════
wb = openpyxl.load_workbook("report.xlsx")
ws = wb["SQL vs SPARQL"]

# Find last data row (before KEY INSIGHT)
last_row = ws.max_row
for r in range(1, ws.max_row + 1):
    if ws.cell(r, 1).value == "KEY INSIGHT":
        last_row = r - 1
        break

ws.insert_rows(last_row + 1)
r = last_row + 1
ws.cell(r,1).value = "QSQL8"
ws.cell(r,2).value = "Total reimplementation family size (any depth, no anchor)"
ws.cell(r,3).value = 7
ws.cell(r,4).value = "861 games returned; top: CATAN 33, Love Letter 25, TtR 21"
ws.cell(r,5).value = 10
ws.cell(r,6).value = "bgg:reimplements+ counts full descendant tree in one token; identical structure to depth-1 query plus one character"
ws.cell(r,7).value = 12
ws.cell(r,8).value = "861 games returned; identical results"
ws.cell(r,9).value = 10
ws.cell(r,10).value = "WITH RECURSIVE needed to compute transitive closure before GROUP BY; requires full query restructure vs SPARQL's single + token"

# Update KEY INSIGHT text to mention QSQL8
for r2 in range(last_row+2, ws.max_row+1):
    v = ws.cell(r2,1).value
    if v and "QSQL7 also shows" in str(v):
        ws.cell(r2,1).value = (
            str(v) + "\n\nQSQL8 (anchorless family size) is the clearest one-token demonstration: "
            "adding + to bgg:reimplements gives full transitive family counts across all 861 games "
            "in the database. SQL needs a WITH RECURSIVE CTE to do the same. Both return identical results."
        )
        break

wb.save("report.xlsx")
print("report.xlsx updated")

# ════════════════════════════════════════════════════════════════════
#  4. UPDATE bgg_sparql_compare.py
# ════════════════════════════════════════════════════════════════════
QSQL8_SPARQL_ENTRY = '''    (
        "QSQL8",
        "Which games have the largest total reimplementation family — "
        "counting all games that have reimplemented them, directly or through any "
        "number of reimplementation generations? Return every game that has been "
        "reimplemented at least once, with its full family size, ordered largest first."
    ),
]'''

txt = open("bgg_sparql_compare.py", encoding="utf-8").read()
txt = txt.replace(
    "    (\n        \"QSQL7\",",
    "    (\n        \"QSQL7\","
)
# Insert QSQL8 just before the closing ] of QUESTIONS
txt = txt.replace(
    "    (\n        \"QSQL7\",\n        \"Show me all games",
    "    (\n        \"QSQL7\",\n        \"Show me all games"  # no-op placeholder
)
# Find the end of the QUESTIONS list and insert before it
marker = ")\n]\n\n# ── Helpers"
if marker in txt:
    txt = txt.replace(marker,
        ")\n" + QSQL8_SPARQL_ENTRY.replace("]\n","") + "\n\n# ── Helpers")
    open("bgg_sparql_compare.py","w",encoding="utf-8").write(txt)
    print("bgg_sparql_compare.py updated")
else:
    print("WARNING: could not find insertion point in bgg_sparql_compare.py — skipping")

# ════════════════════════════════════════════════════════════════════
#  5. UPDATE bgg_sql_compare.py
# ════════════════════════════════════════════════════════════════════
txt2 = open("bgg_sql_compare.py", encoding="utf-8").read()
marker2 = ")\n]\n\n# ── Helpers"
if marker2 in txt2:
    txt2 = txt2.replace(marker2,
        ")\n" + QSQL8_SPARQL_ENTRY.replace("]\n","") + "\n\n# ── Helpers")
    open("bgg_sql_compare.py","w",encoding="utf-8").write(txt2)
    print("bgg_sql_compare.py updated")
else:
    # Try alternative marker
    alt = ")\n]\n\n\n# ── Helpers"
    if alt in txt2:
        txt2 = txt2.replace(alt,
            ")\n" + QSQL8_SPARQL_ENTRY.replace("]\n","") + "\n\n\n# ── Helpers")
        open("bgg_sql_compare.py","w",encoding="utf-8").write(txt2)
        print("bgg_sql_compare.py updated (alt marker)")
    else:
        print("WARNING: could not find insertion point in bgg_sql_compare.py — skipping")

prs.save("bgg_agents_slides.pptx")
print(f"\nDone. Total slides: {len(prs.slides)}")
