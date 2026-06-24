"""Add QSQL4 deep-dive slide to bgg_agents_slides.pptx."""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def box(slide, text, left, top, w, h,
        fill=None, font_size=9, bold=False, color="1A1A2E",
        align=PP_ALIGN.LEFT, italic=False, font_name=None):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    if font_name: run.font.name = font_name
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    return shape

def code_block(slide, code_text, left, top, w, h, font_size=8):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("1C2333")
    first = True
    for line in code_text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Consolas"
        run.font.color.rgb = rgb("D4D4D4")
    return shape

NAVY="1A1A2E"; SPARQL_BLUE="216BC2"; SQL_SLATE="37474F"
WIN_GREEN="1E8C69"; BUG_RED="C0392B"; WHITE="FFFFFF"; ORANGE="D46A08"

prs = Presentation("bgg_agents_slides.pptx")
layout = prs.slide_layouts[6]
s = prs.slides.add_slide(layout)

# Title
box(s, "  QSQL4 — Game Neighbors: Both Queries Correct, Results Uninformative",
    0, 0, 13.33, 0.72, fill=NAVY, font_size=19, bold=True, color=WHITE)

box(s,
    "Who are Susanne Vejdemo's closest game neighbors? "
    "Find the players who share the most games with her and order by overlap descending.",
    0.15, 0.80, 13.00, 0.30, font_size=9.5, italic=True, color="444444")

# ── Queries side by side ──────────────────────────────────────────
SPARQL_Q4 = """\
PREFIX bgg: <https://raw.githubusercontent.com/susvej/bg_ontology/>
PREFIX svj: <https://vejdemo.se/boardgames#>

SELECT ?otherName (COUNT(DISTINCT ?game) AS ?sharedGames) WHERE {
  svj:SusanneVejdemo bgg:hasOwnershipOf ?game .
  ?other bgg:hasOwnershipOf ?game .
  ?other bgg:hasName ?otherName .
  FILTER(?other != svj:SusanneVejdemo)
}
GROUP BY ?other ?otherName
ORDER BY DESC(?sharedGames)"""

SQL_Q4 = """\
SELECT p2.label AS player_name,
       COUNT(DISTINCT po2.game_id) AS shared_games
FROM players p1
JOIN player_owns po1 ON po1.player_id = p1.id
JOIN player_owns po2 ON po2.game_id = po1.game_id
  AND po2.player_id != p1.id
JOIN players p2 ON p2.id = po2.player_id
WHERE p1.label = 'Susanne Vejdemo'
GROUP BY p2.id, p2.label
ORDER BY shared_games DESC
LIMIT 20;"""

box(s, "  SPARQL  (7 lines)", 0.15, 1.18, 6.30, 0.34,
    fill=SPARQL_BLUE, font_size=12, bold=True, color=WHITE)
code_block(s, SPARQL_Q4, 0.15, 1.56, 6.30, 2.60)

box(s, "  SQL  (9 lines)", 6.85, 1.18, 6.30, 0.34,
    fill=SQL_SLATE, font_size=12, bold=True, color=WHITE)
code_block(s, SQL_Q4, 6.85, 1.56, 6.30, 2.60)

# Score badges under each
box(s, "Score: 9/10", 0.15, 4.22, 1.80, 0.30,
    fill=ORANGE, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, "Returned: 41 neighbors (no LIMIT)", 2.05, 4.22, 4.40, 0.30,
    fill="FFF3CD", font_size=8.5, color="5D4037")
box(s, "Score: 9/10", 6.85, 4.22, 1.80, 0.30,
    fill=ORANGE, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, "Returned: 20 neighbors (added LIMIT 20 unprompted)", 8.75, 4.22, 4.40, 0.30,
    fill="FFF3CD", font_size=8.5, color="5D4037")

# ── Results table ─────────────────────────────────────────────────
box(s, "Actual results — top 7 of 41 players who share at least 1 game with Susanne:",
    0.15, 4.62, 13.00, 0.26, font_size=9, bold=True, color="333333")

RESULTS = [
    ("Barbara Lewis",    "2 shared"),
    ("Dorothy Collins",  "1 shared"),
    ("Gertrude Young",   "1 shared"),
    ("Alice Parker",     "1 shared"),
    ("+ 37 more",        "1 shared"),
    ("158 players",      "0 shared"),
    ("(not returned)",   ""),
]
labels = ["1st", "2nd", "3rd", "4th", "...", "", ""]
fills  = ["2C5F2E","37474F","37474F","37474F","37474F","555555","555555"]
for j, ((name, count), lbl, fill) in enumerate(zip(RESULTS, labels, fills)):
    col = j % 4
    row = j // 4
    txt = f"  {lbl:3}  {name:<26}  {count}"
    box(s, txt,
        0.15 + col*3.28, 4.94 + row*0.30, 3.20, 0.28,
        fill=fill, font_size=8, color="D4D4D4", font_name="Consolas")

# ── Root cause callout ────────────────────────────────────────────
box(s, "Root cause: data sparsity",
    0.15, 5.64, 4.60, 0.32, fill="C0392B", font_size=10, bold=True, color=WHITE)
box(s,
    "200 players each own ~30 games drawn from a pool of 21,000. "
    "Probability of any two players sharing a game: 30/21,000 x 30 = 0.04 expected overlaps per pair. "
    "Result: 41 players share 1 game, 1 player shares 2. "
    "The maximum possible signal in this data is almost zero.",
    0.15, 6.00, 6.20, 0.54, fill="FDECEA", font_size=8.5, color="5D0000")

# ── What both agents did wrong ────────────────────────────────────
box(s, "What both agents got wrong",
    6.85, 5.64, 6.30, 0.32, fill="37474F", font_size=10, bold=True, color=WHITE)

lines = [
    ("Query logic", "Both correct — valid GROUP BY + COUNT join", WIN_GREEN),
    ("SPARQL",      "No LIMIT: returned all 41 neighbors, not just the top", ORANGE),
    ("SQL",         "Added LIMIT 20 unprompted — arbitrary cutoff", ORANGE),
    ("Both",        "Neither flagged data sparsity or explained max=2 is an artefact", BUG_RED),
    ("Fix needed",  "Clustered players with taste overlap would expose real structure", SPARQL_BLUE),
]
for k, (label, text, pill_fill) in enumerate(lines):
    y_pos = 6.00 + k * 0.28
    box(s, f" {label} ", 6.85, y_pos, 1.42, 0.24,
        fill=pill_fill, font_size=8, bold=True, color=WHITE)
    box(s, text, 8.35, y_pos, 4.80, 0.24,
        font_size=8, color="333333")

prs.save("bgg_agents_slides.pptx")
print(f"QSQL4 deep-dive slide added. Total slides: {len(prs.slides)}")
